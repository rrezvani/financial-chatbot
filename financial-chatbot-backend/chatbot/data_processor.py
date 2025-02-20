import pandas as pd
import fitz  # PyMuPDF for PDF processing
from pptx import Presentation
from sentence_transformers import SentenceTransformer
import faiss
import networkx as nx
import json
from pathlib import Path
import numpy as np
import PyPDF2
import pptx
import os
from . import config
import re
from datetime import datetime

class DataProcessor:
    def __init__(self):
        # Initialize sentence transformer model
        self.model = SentenceTransformer(config.EMBEDDING_MODEL)
        
        # Store processed text chunks
        self.documents = []
        
        # Vector store for semantic search
        self.vector_store = None
        
        # Knowledge graph for relationships
        self.graph = nx.Graph()
        
        # Initialize dictionaries for tax data
        self.tax_rates = {}
        self.tax_rules = {}
        self.tax_examples = {}
        
        # Define relationship types
        self.relationship_types = {
            'TAX_RATE': 'has_tax_rate',
            'DEDUCTION': 'has_deduction',
            'APPLIES_TO': 'applies_to',
            'RELATED_TO': 'related_to',
            'VALID_IN': 'valid_in'
        }
        
        # Add entity types mapping
        self.entity_types = {
            'non-profit': 'Non-Profit',
            'nonprofit': 'Non-Profit',
            'corporation': 'Corporation',
            'corporate': 'Corporation',
            'individual': 'Individual',
            'person': 'Individual',
            'trust': 'Trust',
            'partnership': 'Partnership'
        }
        
        # Process our tax documents
        self.load_tax_documents()

    def load_tax_documents(self):
        """Load and process all tax documents once"""
        try:
            # Load CSV data
            print(f"Attempting to load CSV from: {config.TAX_DOCUMENTS['tax_data']}")
            tax_rates = pd.read_csv(config.TAX_DOCUMENTS['tax_data'])
            if tax_rates is not None:
                print("CSV loaded successfully, processing tax rates...")
                self.process_tax_rates(tax_rates)
            
            # Load main tax code PDF
            try:
                print(f"Attempting to load main tax code from: {config.TAX_DOCUMENTS['tax_code']}")
                with open(config.TAX_DOCUMENTS['tax_code'], 'rb') as pdf_file:
                    pdf_reader = PyPDF2.PdfReader(pdf_file)
                    self.process_tax_guidelines(pdf_reader, "tax_code")
            except FileNotFoundError as e:
                print(f"Warning: Tax code PDF not found at {config.TAX_DOCUMENTS['tax_code']}")
                print(f"Error details: {str(e)}")
            
            # Load tax instructions PDF
            try:
                print(f"Attempting to load tax instructions from: {config.TAX_DOCUMENTS['tax_instructions']}")
                with open(config.TAX_DOCUMENTS['tax_instructions'], 'rb') as pdf_file:
                    pdf_reader = PyPDF2.PdfReader(pdf_file)
                    self.process_tax_guidelines(pdf_reader, "instructions")
            except FileNotFoundError as e:
                print(f"Warning: Tax instructions PDF not found at {config.TAX_DOCUMENTS['tax_instructions']}")
                print(f"Error details: {str(e)}")
            
            # Load PPT data
            ppt_path = config.TAX_DOCUMENTS['tax_presentation']
            print(f"Checking if PPT exists at: {ppt_path}")
            if os.path.exists(ppt_path):
                print("PPT file found, attempting to load...")
                try:
                    prs = pptx.Presentation(ppt_path)
                    self.process_tax_examples(prs)
                except Exception as e:
                    print(f"Error loading PPT: {str(e)}")
                    print(f"Error type: {type(e)}")
            else:
                print(f"PPT file not found at: {ppt_path}")
            
            print("Building search index...")
            self.build_search_index()
            
            # Add this line to build the knowledge graph
            print("Building knowledge graph...")
            self.build_knowledge_graph()
        except Exception as e:
            print(f"Error loading tax documents: {str(e)}")
            print(f"Error type: {type(e)}")
            import traceback
            print(f"Full traceback: {traceback.format_exc()}")

    def process_tax_rates(self, df):
        """Process tax rates from CSV"""
        try:
            print("Available columns:", df.columns.tolist())
            
            # Sort by income to ensure proper bracketing
            df = df.sort_values('Income')
            
            for _, row in df.iterrows():
                try:
                    # Clean and convert numeric values
                    income = float(str(row['Income']).replace('$', '').replace(',', ''))
                    deductions = float(str(row['Deductions']).replace('$', '').replace(',', ''))
                    taxable_income = float(str(row['Taxable Income']).replace('$', '').replace(',', ''))
                    tax_owed = float(str(row['Tax Owed']).replace('$', '').replace(',', ''))
                    
                    # Fix tax rate conversion - store as actual percentage value
                    tax_rate_str = str(row['Tax Rate']).strip().rstrip('%')
                    tax_rate = float(tax_rate_str) * 100  # Convert decimal to percentage
                    
                    bracket_id = f"bracket_{income}"
                    
                    self.tax_rates[bracket_id] = {
                        'range': f"${income:,.2f}",
                        'rate': tax_rate,  # Now storing as actual percentage (e.g., 24.55)
                        'taxpayer_type': str(row['Taxpayer Type']),
                        'tax_year': str(row['Tax Year']),
                        'transaction_date': pd.to_datetime(row['Transaction Date']).strftime('%Y-%m-%d'),
                        'income_source': str(row['Income Source']),
                        'deduction_type': str(row['Deduction Type']),
                        'state': str(row['State']),
                        'deductions': f"${deductions:,.2f}",
                        'taxable_income': f"${taxable_income:,.2f}",
                        'tax_owed': f"${tax_owed:,.2f}",
                        'conditions': f"Type: {row['Taxpayer Type']}, Year: {row['Tax Year']}, State: {row['State']}"
                    }
                    
                except ValueError as e:
                    print(f"Warning: Error processing row {_}: {e}")
                    continue

        except Exception as e:
            print(f"Error processing tax rates: {str(e)}")
            import traceback
            print(f"Full traceback: {traceback.format_exc()}")

    def process_tax_guidelines(self, pdf_reader, source):
        """Process tax guidelines from PDF"""
        # Extract text directly from the PdfReader object
        chunks = []
        
        # Get total number of pages
        total_pages = len(pdf_reader.pages)
        print(f"PDF ({source}) has {total_pages} pages. Processing first 20 pages...")
        
        # Only process first 20 pages - most tax guidelines have important info up front
        for page_num in range(min(20, total_pages)):
            if page_num % 5 == 0:  # Progress update more frequently
                print(f"Processing page {page_num}...")
                
            try:
                text = pdf_reader.pages[page_num].extract_text()
                
                # More aggressive filtering
                if len(text.strip()) < 100:  # Skip very short pages
                    continue
                    
                # Only keep paragraphs that mention relevant tax terms
                relevant_terms = ['tax', 'income', 'deduction', 'rate', 'bracket']
                paragraphs = text.split('\n\n')
                relevant_paragraphs = [
                    p for p in paragraphs 
                    if any(term in p.lower() for term in relevant_terms)
                ]
                
                # Join relevant paragraphs and add to chunks
                if relevant_paragraphs:
                    chunks.append(' '.join(relevant_paragraphs))
                    
            except Exception as e:
                print(f"Error processing page {page_num}: {str(e)}")
                continue

        print(f"Processed {len(chunks)} relevant chunks from the PDF")
        print("Building relationships for relevant chunks only...")

        # Only process chunks that are highly relevant
        for i, chunk in enumerate(chunks):
            rule_id = f"rule_{source}_{i}"  # Add source to ID
            self.tax_rules[rule_id] = {
                'text': chunk,
                'source': source  # Use the provided source
            }
            self.documents.append(chunk)
            
            # Add to graph without checking relationships
            # We'll rely on vector search instead
            self.graph.add_node(rule_id, 
                              type='tax_rule',
                              data=self.tax_rules[rule_id])

        print("Finished processing PDF chunks")

    def process_tax_examples(self, ppt):
        """Process tax examples from PPT"""
        print("Starting to process PPT...")
        try:
            chunks = []
            print(f"Processing {len(ppt.slides)} slides...")
            
            for slide in ppt.slides:
                slide_text = ""
                
                # Get text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text += shape.text + " "
                        
                # Get text from notes
                if slide.has_notes_slide:
                    notes = slide.notes_slide.notes_text_frame.text
                    slide_text += " [Notes: " + notes + "]"
                    
                if slide_text.strip():
                    chunks.append(slide_text.strip())
            
            print(f"Extracted {len(chunks)} chunks from PPT")
            
            for i, chunk in enumerate(chunks):
                if i % 10 == 0:  # Progress update
                    print(f"Processing chunk {i}...")
                    
                example_id = f"example_{i}"
                self.tax_examples[example_id] = {
                    'text': chunk,
                    'source': 'examples'
                }
                self.documents.append(chunk)  # Add to documents for search
                
                # Add to graph
                self.graph.add_node(example_id,
                                  type='example',
                                  data=self.tax_examples[example_id])
                    
            print("Finished processing PPT")
            
        except Exception as e:
            print(f"Error processing PPT: {str(e)}")
            print(f"Error type: {type(e)}")
            import traceback
            print(f"Full traceback: {traceback.format_exc()}")

    def build_search_index(self):
        """Build vector store and knowledge graph"""
        try:
            print(f"Building search index for {len(self.documents)} documents...")
            
            if not self.documents:
                print("Warning: No documents to index")
                return
            
            # Process in batches to avoid memory issues
            batch_size = 32
            all_embeddings = []
            
            for i in range(0, len(self.documents), batch_size):
                batch = self.documents[i:i + batch_size]
                print(f"Processing batch {i//batch_size + 1}/{(len(self.documents) + batch_size - 1)//batch_size}")
                
                # Create embeddings for current batch
                batch_embeddings = self.model.encode(batch)
                all_embeddings.append(batch_embeddings)
                
            # Concatenate all embeddings
            print("Concatenating embeddings...")
            embeddings = np.vstack(all_embeddings)
            
            # Build FAISS index
            print("Building FAISS index...")
            dimension = embeddings.shape[1]
            self.vector_store = faiss.IndexFlatL2(dimension)
            self.vector_store.add(embeddings.astype('float32'))
            
            print("Search index built successfully!")
            
        except Exception as e:
            print(f"Error building search index: {str(e)}")
            print(f"Error type: {type(e)}")
            import traceback
            print(f"Full traceback: {traceback.format_exc()}")

    def vector_search(self, query: str) -> list:
        """Perform vector-based search"""
        query_vector = self.model.encode([query])
        D, I = self.vector_store.search(query_vector.astype('float32'), k=config.TOP_K_RESULTS)
        
        matches = []
        for i, score in zip(I[0], D[0]):
            if i < len(self.documents):
                text = self.documents[i]
                # Format any percentage matches in the text
                percentage_matches = re.findall(r'(\d+\.?\d*)%', text)
                for match in percentage_matches:
                    text = text.replace(f"{match}%", f"{float(match):.2f}%")
                
                matches.append({
                    'text': text,
                    'score': float(score)
                })
        return matches

    def search(self, query: str):
        """Combined vector and graph-based search"""
        try:
            query = query.strip('- ').strip()
            query_lower = query.lower()
            
            # Handle income source queries
            if 'income' in query_lower and any(source.lower() in query_lower for source in 
                ['business income', 'capital gains', 'rental', 'salary', 'royalties', 'investment']):
                return self.calculate_average_by_source(query)
            
            # Handle deduction type queries
            if any(deduction.lower() in query_lower for deduction in 
                ['education expenses', 'medical expenses', 'mortgage interest', 'business expenses', 'charitable contributions']):
                return self.calculate_average_by_deduction(query)
            
            # Handle average queries before income check
            if 'average' in query_lower or 'mean' in query_lower:
                return self.calculate_average(query)
            
            # Handle income-based queries
            income_match = re.search(r'-?\$?([\d,]+(?:\.\d{2})?)', query)
            if income_match:
                # Check for negative sign in original query
                if '-' in query:
                    return {
                        'direct_matches': ['Invalid query: Please provide a positive income amount.'],
                        'enhanced_results': []
                    }
                
                query_income = float(income_match.group(1).replace(',', ''))
                if query_income > 999999999:
                    return {
                        'direct_matches': ['Invalid query: Please provide an income amount less than $1,000,000,000.'],
                        'enhanced_results': []
                    }
                return self.get_income_based_rates(query_income)

            # Edge case validation
            if any(term in query_lower for term in ['billion', 'trillion', 'million']):
                return {
                    'direct_matches': ['Invalid query: Please provide a valid income amount between $0 and $999,999,999.'],
                    'enhanced_results': []
                }

            # Handle comparison queries
            if 'compare' in query_lower or 'difference between' in query_lower:
                # State comparison
                if 'state' in query_lower or any(f" {state.lower()} " in f" {query_lower} " 
                    for state in ['TX', 'CA', 'NY', 'FL', 'PA', 'IL', 'OH', 'GA', 'NC', 'MI']):
                    states = re.findall(r'\b([A-Z]{2})\b', query.upper())
                    if len(states) == 2:
                        return self.compare_states(states[0], states[1], query)
                
                # Organization type comparison
                if any(org in query_lower for org in ['corporation', 'non-profit', 'trust', 'individual', 'partnership']):
                    return self.compare_organizations(query)

            # Remove highest state tax rate handling and use default response
            if 'highest' in query_lower and 'state' in query_lower:
                return {
                    'direct_matches': ['I apologize, but I cannot determine the highest state tax rate at this time.'],
                    'enhanced_results': []
                }

            # Default to vector search
            return self.handle_general_query(query)

        except Exception as e:
            print(f"Error in search: {str(e)}")
            return {
                'direct_matches': ['Sorry, I encountered an error while searching.'],
                'enhanced_results': []
            }

    def compare_organizations(self, query: str) -> dict:
        """Compare tax rates between different organization types"""
        org_types = {
            'corporation': ['corporation', 'corporate'],
            'non-profit': ['non-profit', 'nonprofit'],
            'individual': ['individual', 'person', 'personal'],
            'trust': ['trust'],
            'partnership': ['partnership']
        }
        
        found_types = []
        for org_type, keywords in org_types.items():
            if any(keyword in query.lower() for keyword in keywords):
                found_types.append(org_type)
        
        if len(found_types) < 2:
            return {
                'direct_matches': ['Please specify two organization types to compare.'],
                'enhanced_results': []
            }
        
        results = {}
        for org_type in found_types:
            rates = []
            for tax_info in self.tax_rates.values():
                if org_type.title() in tax_info['conditions']:
                    rates.append(tax_info['rate'])
            
            if rates:
                avg_rate = sum(rates) / len(rates)
                results[org_type] = {
                    'average_rate': avg_rate,
                    'sample_size': len(rates)
                }
        
        if not results:
            return {
                'direct_matches': ['No comparable data found for the specified organization types.'],
                'enhanced_results': []
            }
        
        response = ['Comparison of average tax rates:']
        for org_type, data in results.items():
            response.append(
                f"{org_type.title()}: {self.format_rate(data['average_rate'])} "
                f"(based on {data['sample_size']} records)"
            )
        
        return {
            'direct_matches': response,
            'enhanced_results': []
        }

    def calculate_average(self, query: str) -> dict:
        """Calculate average tax rates based on query criteria"""
        query_lower = query.lower()
        
        # Extract organization type
        org_type = None
        for key, value in self.entity_types.items():
            if key in query_lower:
                org_type = value
                break
        
        # Extract year
        year = None
        year_match = re.search(r'\b(19|20)\d{2}\b', query)
        if year_match:
            year = year_match.group(0)
        
        matching_rates = []
        for tax_info in self.tax_rates.values():
            matches_criteria = True
            
            # Check organization type
            if org_type and org_type not in tax_info['taxpayer_type']:
                matches_criteria = False
            
            # Check year
            if year and str(year) not in tax_info['tax_year']:
                matches_criteria = False
            
            if matches_criteria:
                matching_rates.append(tax_info['rate'])
        
        if not matching_rates:
            return {
                'direct_matches': [f'No tax rates found for {org_type or "organizations"} in {year or "any year"}.'],
                'enhanced_results': []
            }
        
        avg_rate = sum(matching_rates) / len(matching_rates)
        criteria = []
        if org_type:
            criteria.append(org_type)
        if year:
            criteria.append(f"in {year}")
        
        criteria_str = ' '.join(criteria) if criteria else 'all records'
        
        return {
            'direct_matches': [
                f"The average tax rate for {criteria_str} is {self.format_rate(avg_rate)}",
                f"This is based on {len(matching_rates)} records"
            ],
            'enhanced_results': []
        }

    def calculate_average_by_source(self, query: str) -> dict:
        """Calculate average tax rate for a specific income source"""
        query_lower = query.lower()
        
        # Map common variations to standardized names
        source_mapping = {
            'business': 'Business Income',
            'capital': 'Capital Gains',
            'rental': 'Rental',
            'salary': 'Salary',
            'royalties': 'Royalties',
            'investment': 'Investment'
        }
        
        income_source = None
        for key, value in source_mapping.items():
            if key in query_lower:
                income_source = value
                break
        
        if not income_source:
            return {
                'direct_matches': ['Please specify a valid income source.'],
                'enhanced_results': []
            }
        
        matching_rates = []
        for tax_info in self.tax_rates.values():
            if tax_info['income_source'] == income_source:
                matching_rates.append(tax_info['rate'])
        
        if not matching_rates:
            return {
                'direct_matches': [f'No tax rates found for {income_source}.'],
                'enhanced_results': []
            }
        
        avg_rate = sum(matching_rates) / len(matching_rates)
        return {
            'direct_matches': [
                f"The average tax rate for {income_source} is {self.format_rate(avg_rate)}",
                f"This is based on {len(matching_rates)} records"
            ],
            'enhanced_results': []
        }

    def calculate_average_by_deduction(self, query: str) -> dict:
        """Calculate average tax rate for a specific deduction type"""
        query_lower = query.lower()
        
        # Map common variations to standardized names
        deduction_mapping = {
            'education': 'Education Expenses',
            'medical': 'Medical Expenses',
            'mortgage': 'Mortgage Interest',
            'business expense': 'Business Expenses',
            'charitable': 'Charitable Contributions'
        }
        
        deduction_type = None
        for key, value in deduction_mapping.items():
            if key in query_lower:
                deduction_type = value
                break
        
        if not deduction_type:
            return {
                'direct_matches': ['Please specify a valid deduction type.'],
                'enhanced_results': []
            }
        
        matching_rates = []
        for tax_info in self.tax_rates.values():
            if tax_info['deduction_type'] == deduction_type:
                matching_rates.append(tax_info['rate'])
        
        if not matching_rates:
            return {
                'direct_matches': [f'No tax rates found for {deduction_type}.'],
                'enhanced_results': []
            }
        
        avg_rate = sum(matching_rates) / len(matching_rates)
        return {
            'direct_matches': [
                f"The average tax rate for those with {deduction_type} is {self.format_rate(avg_rate)}",
                f"This is based on {len(matching_rates)} records"
            ],
            'enhanced_results': []
        }

    def extract_text_from_pdf(self, pdf_path):
        """Extract text from PDF and split into chunks"""
        chunks = []
        with open(pdf_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            current_chunk = ""
            
            for page in pdf_reader.pages:
                text = page.extract_text()
                
                # Split into sentences (rough approximation)
                sentences = text.replace('\n', ' ').split('. ')
                
                for sentence in sentences:
                    if len(current_chunk) + len(sentence) < config.MAX_CHUNK_SIZE:
                        current_chunk += sentence + '. '
                    else:
                        chunks.append(current_chunk)
                        current_chunk = sentence + '. '
                        
            if current_chunk:
                chunks.append(current_chunk)
                
        return chunks

    def format_rate(self, rate: float) -> str:
        """Format tax rate as percentage string"""
        return f"{rate:.2f}%"  # rate is already a percentage value

    def is_related(self, query_income: float, tax_info: dict) -> bool:
        """Check if a tax bracket is related to the query income"""
        try:
            bracket_income = float(tax_info['range'].replace('$', '').replace(',', ''))
            # Look for brackets within 5% above OR below the query income
            percentage_diff = abs(query_income - bracket_income) / query_income
            return percentage_diff < 0.05
        except (ValueError, TypeError):
            return False

    def calculate_similarity(self, text1: str, text2: str) -> float:
        """Calculate cosine similarity between two texts"""
        embeddings = self.model.encode([text1, text2])
        similarity = np.dot(embeddings[0], embeddings[1]) / (
            np.linalg.norm(embeddings[0]) * np.linalg.norm(embeddings[1])
        )
        return similarity

    def build_knowledge_graph(self):
        """Build knowledge graph from processed data"""
        # Add nodes and edges from tax rates
        for bracket_id, tax_info in self.tax_rates.items():
            try:
                # Create unique identifiers
                rate_id = f"rate_{bracket_id}"
                income_id = f"income_{bracket_id}"
                state = tax_info['conditions'].split(', ')[-1].split(': ')[1]
                year = tax_info['conditions'].split(', ')[1].split(': ')[1]
                tax_type = tax_info['conditions'].split(', ')[0].split(': ')[1]

                state_id = f"state_{state}"
                year_id = f"year_{year}"
                type_id = f"type_{tax_type}"

                # Add nodes with attributes
                self.graph.add_node(rate_id, 
                                  type='tax_rate', 
                                  rate=float(tax_info['rate']),  # tax_info['rate'] is already a float
                                  range=float(tax_info['range'].replace('$', '').replace(',', '')))
                
                self.graph.add_node(income_id, type='income_bracket')
                self.graph.add_node(state_id, type='state')
                self.graph.add_node(year_id, type='year')
                self.graph.add_node(type_id, type='taxpayer_type')

                # Add relationships
                self.graph.add_edge(rate_id, income_id, relationship=self.relationship_types['APPLIES_TO'])
                self.graph.add_edge(rate_id, state_id, relationship=self.relationship_types['VALID_IN'])
                self.graph.add_edge(rate_id, year_id, relationship=self.relationship_types['VALID_IN'])
                self.graph.add_edge(rate_id, type_id, relationship=self.relationship_types['APPLIES_TO'])
            except Exception as e:
                print(f"Warning: Error processing bracket {bracket_id}: {e}")
                continue

        # Add relationships from text chunks
        for i, doc in enumerate(self.documents):
            chunk_id = f"chunk_{i}"
            self.graph.add_node(chunk_id, 
                              type='text_chunk',
                              text=doc)  # Fix: doc is the text directly

            # Connect to related tax rates
            for node in self.graph.nodes():
                node_data = self.graph.nodes[node]
                if node_data.get('type') == 'tax_rate':
                    if self._are_related(doc, node_data):  # Fix: pass doc text directly
                        self.graph.add_edge(chunk_id, node, 
                                          relationship=self.relationship_types['RELATED_TO'])

    def _are_related(self, text: str, node_data: dict) -> bool:
        """Check if text chunk is related to a node"""
        if node_data.get('type') == 'tax_rate':
            try:
                # Safely get rate and range with defaults
                rate = node_data.get('rate', 0)
                range_val = node_data.get('range', 0)
                
                # Only proceed if we have both values
                if rate and range_val:
                    rate_str = f"{float(rate)*100:.1f}%"
                    income_str = f"${float(range_val):,.2f}"
                    return rate_str in text or income_str in text
            except (ValueError, TypeError) as e:
                print(f"Warning: Error processing node data: {e}")
                return False
        return False

    def find_related_info(self, query: str, max_results: int = 3) -> list:
        """Find related information using both vector and graph-based approaches"""
        # Get initial matches using vector search
        vector_matches = self.vector_search(query)
        
        # Extract relevant nodes from vector matches
        relevant_nodes = set()
        for match in vector_matches:
            # Find connected nodes in graph
            for node in self.graph.nodes():
                if match['text'] in str(self.graph.nodes[node].get('text', '')):
                    relevant_nodes.add(node)
                    # Add neighbors up to 2 hops away
                    relevant_nodes.update(nx.single_source_shortest_path_length(self.graph, node, cutoff=2).keys())

        # Score and rank results
        results = []
        for node in relevant_nodes:
            node_data = self.graph.nodes[node]
            if node_data.get('type') == 'tax_rate':
                score = self._calculate_relevance_score(query, node, node_data)
                results.append({
                    'type': 'tax_rate',
                    'rate': f"{node_data['rate']*100:.2f}%",
                    'income_range': f"${node_data['range']:,.2f}",
                    'conditions': self._get_node_conditions(node),
                    'score': score
                })

        # Sort by relevance score and return top results
        results.sort(key=lambda x: x['score'], reverse=True)
        return results[:max_results]

    def _calculate_relevance_score(self, query: str, node_id: str, node_data: dict) -> float:
        """Calculate relevance score based on multiple factors"""
        score = 0.0
        
        # Base score from vector similarity
        if node_data.get('text'):
            score += self.calculate_similarity(query, node_data['text'])
        
        # Graph-based factors
        # More connections = more relevant
        score += len(list(self.graph.neighbors(node_id))) * 0.1
        
        # Prefer more recent information
        if 'year' in str(node_id):
            try:
                year = int(str(node_id).split('_')[1])
                current_year = datetime.now().year
                score += 1.0 / (current_year - year + 1)
            except:
                pass
        
        return score

    def _get_node_conditions(self, node_id: str) -> str:
        """Get conditions associated with a node from its relationships"""
        conditions = []
        for neighbor in self.graph.neighbors(node_id):
            neighbor_data = self.graph.nodes[neighbor]
            edge_data = self.graph.edges[node_id, neighbor]
            
            if neighbor_data.get('type') in ['state', 'year', 'taxpayer_type']:
                value = neighbor.split('_')[1]
                conditions.append(f"{neighbor_data['type'].title()}: {value}")
                
        return ', '.join(conditions) 

    def compare_states(self, state1: str, state2: str, query: str) -> dict:
        """Compare tax rates between two states"""
        try:
            # Extract income if present
            income_match = re.search(r'\$?([\d,]+(?:\.\d{2})?)', query)
            income = float(income_match.group(1).replace(',', '')) if income_match else None
            
            state1_rates = []
            state2_rates = []
            
            for tax_info in self.tax_rates.values():
                bracket_income = float(tax_info['range'].replace('$', '').replace(',', ''))
                
                # If income specified, only compare nearby brackets
                if income and abs(income - bracket_income) / income > 0.05:
                    continue
                    
                if state1 in tax_info['conditions']:
                    state1_rates.append({
                        'rate': tax_info['rate'],
                        'income': bracket_income,
                        'type': tax_info['taxpayer_type']
                    })
                elif state2 in tax_info['conditions']:
                    state2_rates.append({
                        'rate': tax_info['rate'],
                        'income': bracket_income,
                        'type': tax_info['taxpayer_type']
                    })
            
            if not state1_rates or not state2_rates:
                return {
                    'direct_matches': [f'No comparable tax rates found for {state1} and {state2}.'],
                    'enhanced_results': []
                }
            
            # Calculate averages
            avg1 = sum(r['rate'] for r in state1_rates) / len(state1_rates)
            avg2 = sum(r['rate'] for r in state2_rates) / len(state2_rates)
            
            return {
                'direct_matches': [
                    f"Comparing tax rates between {state1} and {state2}:",
                    f"{state1} average rate: {self.format_rate(avg1)}",
                    f"{state2} average rate: {self.format_rate(avg2)}",
                    f"Based on {len(state1_rates)} records for {state1} and {len(state2_rates)} records for {state2}"
                ],
                'enhanced_results': []
            }
        except Exception as e:
            print(f"Error in state comparison: {str(e)}")
            return {
                'direct_matches': ['Error comparing states. Please try a different query.'],
                'enhanced_results': []
            } 

    def get_income_based_rates(self, query_income: float, query: str = '') -> dict:
        """Get tax rates for a specific income level with optional filters"""
        query_lower = query.lower()
        
        # Extract state if present
        state = None
        state_match = re.findall(r'\b([A-Z]{2})\b', query.upper())
        if state_match:
            state = state_match[0]
        
        # Extract organization type if present
        org_type = None
        for key, value in self.entity_types.items():
            if key in query_lower:
                org_type = value
                break
        
        matching_brackets = []
        for tax_info in self.tax_rates.values():
            # Check income range
            if not self.is_related(query_income, tax_info):
                continue
            
            # Check state if specified
            if state and state not in tax_info['state']:
                continue
            
            # Check organization type if specified
            if org_type and org_type not in tax_info['taxpayer_type']:
                continue
            
            matching_brackets.append({
                'type': 'tax_rate',
                'income_range': tax_info['range'],
                'rate': self.format_rate(tax_info['rate']),
                'conditions': tax_info['conditions'],
                'details': {
                    'income_source': tax_info['income_source'],
                    'deduction_type': tax_info['deduction_type'],
                    'deductions': tax_info['deductions'],
                    'taxable_income': tax_info['taxable_income'],
                    'tax_owed': tax_info['tax_owed'],
                    'transaction_date': tax_info['transaction_date']
                }
            })
        
        # Sort by how close the income is to the query income
        matching_brackets.sort(key=lambda x: abs(float(x['income_range'].replace('$', '').replace(',', '')) - query_income))
        
        if matching_brackets:
            return {
                'direct_matches': ['Based on the income provided, here are the applicable tax rates:'],
                'enhanced_results': matching_brackets[:3]
            }
        else:
            return {
                'direct_matches': ['No tax rates found matching all specified criteria.'],
                'enhanced_results': []
            }

    def handle_general_query(self, query: str) -> dict:
        """Handle queries that don't match specific patterns"""
        vector_matches = self.vector_search(query)
        if not vector_matches:
            return {
                'direct_matches': ['I could not find specific information to answer your question.'],
                'enhanced_results': []
            }
        
        response = []
        for match in vector_matches[:2]:  # Take top 2 matches
            content = match['text'].replace('\n', ' ').strip()
            if len(content) > 50:  # Only include substantial matches
                response.append(content)
        
        return {
            'direct_matches': response if response else ['No relevant information found.'],
            'enhanced_results': []
        } 