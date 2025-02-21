import pandas as pd
import fitz  # PyMuPDF for PDF processing
from pptx import Presentation
from sentence_transformers import SentenceTransformer
import faiss
import networkx as nx
import json
from pathlib import Path
import numpy as np
import PyPDF2
import pptx
import os
from . import config
import re
from datetime import datetime
import traceback  # Add this too for better error tracking

class DataProcessor:
    def __init__(self):
        # Initialize sentence transformer model
        self.model = SentenceTransformer(config.EMBEDDING_MODEL)
        
        # Store processed text chunks
        self.documents = []
        
        # Vector store for semantic search
        self.vector_store = None
        
        # Knowledge graph for relationships
        self.graph = nx.Graph()
        
        # Initialize dictionaries for tax data
        self.tax_rates = {}
        self.tax_rules = {}
        self.tax_examples = {}
        
        # Define relationship types
        self.relationship_types = {
            'TAX_RATE': 'has_tax_rate',
            'DEDUCTION': 'has_deduction',
            'APPLIES_TO': 'applies_to',
            'RELATED_TO': 'related_to',
            'VALID_IN': 'valid_in'
        }
        
        # Add entity types mapping
        self.entity_types = {
            'non-profit': 'Non-Profit',
            'nonprofit': 'Non-Profit',
            'corporation': 'Corporation',
            'corporate': 'Corporation',
            'individual': 'Individual',
            'person': 'Individual',
            'trust': 'Trust',
            'partnership': 'Partnership'
        }
        
        # Added for new functionality
        self.graph_data = {}
        self.table_data = {}
        self.graph_relationships = []
        
        # Process our tax documents
        self.load_tax_documents()

    def load_tax_documents(self):
        """Load and process all tax documents once"""
        try:
            # Load CSV data
            print(f"Attempting to load CSV from: {config.TAX_DOCUMENTS['tax_data']}")
            tax_rates = pd.read_csv(config.TAX_DOCUMENTS['tax_data'])
            if tax_rates is not None:
                print("CSV loaded successfully, processing tax rates...")
                self.process_tax_rates(tax_rates)
            
            # Load main tax code PDF
            try:
                print(f"Attempting to load main tax code from: {config.TAX_DOCUMENTS['tax_code']}")
                with open(config.TAX_DOCUMENTS['tax_code'], 'rb') as pdf_file:
                    pdf_reader = PyPDF2.PdfReader(pdf_file)
                    self.process_tax_guidelines(pdf_reader, "tax_code")
            except FileNotFoundError as e:
                print(f"Warning: Tax code PDF not found at {config.TAX_DOCUMENTS['tax_code']}")
                print(f"Error details: {str(e)}")
            
            # Load tax instructions PDF
            try:
                print(f"Attempting to load tax instructions from: {config.TAX_DOCUMENTS['tax_instructions']}")
                with open(config.TAX_DOCUMENTS['tax_instructions'], 'rb') as pdf_file:
                    pdf_reader = PyPDF2.PdfReader(pdf_file)
                    self.process_tax_guidelines(pdf_reader, "instructions")
            except FileNotFoundError as e:
                print(f"Warning: Tax instructions PDF not found at {config.TAX_DOCUMENTS['tax_instructions']}")
                print(f"Error details: {str(e)}")
            
            # Load PPT data
            ppt_path = config.TAX_DOCUMENTS['tax_presentation']
            print("\nDEBUG: ===== PPT Loading =====")
            print(f"DEBUG: PPT path: {ppt_path}")
            print(f"DEBUG: Path exists: {os.path.exists(ppt_path)}")
            
            if os.path.exists(ppt_path):
                print("DEBUG: Found PPT file")
                try:
                    print("DEBUG: Starting PPT processing")
                    prs = Presentation(ppt_path)
                    print(f"DEBUG: Loaded presentation with {len(prs.slides)} slides")
                    
                    # Initialize data structures
                    self.graph_data = {}
                    self.table_data = {}
                    
                    # Process the PPT
                    self.process_ppt(prs)
                    
                    print("\nDEBUG: After processing:")
                    print(f"DEBUG: Graph data keys: {list(self.graph_data.keys())}")
                    print(f"DEBUG: Table data keys: {list(self.table_data.keys())}")
                    
                except Exception as e:
                    print(f"DEBUG: Error processing PPT: {str(e)}")
                    traceback.print_exc()
            else:
                print(f"DEBUG: PPT file not found at {ppt_path}")
            
            print("Building search index...")
            self.build_search_index()
            
            # Add this line to build the knowledge graph
            print("Building knowledge graph...")
            self.build_knowledge_graph()
        except Exception as e:
            print(f"Error loading tax documents: {str(e)}")
            print(f"Error type: {type(e)}")
            import traceback
            print(f"Full traceback: {traceback.format_exc()}")

    def process_tax_rates(self, df):
        """Process tax rates from CSV"""
        try:
            print("Available columns:", df.columns.tolist())
            
            # Sort by income to ensure proper bracketing
            df = df.sort_values('Income')
            
            for _, row in df.iterrows():
                try:
                    # Clean and convert numeric values
                    income = float(str(row['Income']).replace('$', '').replace(',', ''))
                    deductions = float(str(row['Deductions']).replace('$', '').replace(',', ''))
                    taxable_income = float(str(row['Taxable Income']).replace('$', '').replace(',', ''))
                    tax_owed = float(str(row['Tax Owed']).replace('$', '').replace(',', ''))
                    
                    # Fix tax rate conversion - store as actual percentage value
                    tax_rate_str = str(row['Tax Rate']).strip().rstrip('%')
                    tax_rate = float(tax_rate_str) * 100  # Convert decimal to percentage
                    
                    bracket_id = f"bracket_{income}"
                    
                    self.tax_rates[bracket_id] = {
                        'range': f"${income:,.2f}",
                        'rate': tax_rate,  # Now storing as actual percentage (e.g., 24.55)
                        'taxpayer_type': str(row['Taxpayer Type']),
                        'tax_year': str(row['Tax Year']),
                        'transaction_date': pd.to_datetime(row['Transaction Date']).strftime('%Y-%m-%d'),
                        'income_source': str(row['Income Source']),
                        'deduction_type': str(row['Deduction Type']),
                        'state': str(row['State']),
                        'deductions': f"${deductions:,.2f}",
                        'taxable_income': f"${taxable_income:,.2f}",
                        'tax_owed': f"${tax_owed:,.2f}",
                        'conditions': f"Type: {row['Taxpayer Type']}, Year: {row['Tax Year']}, State: {row['State']}"
                    }
                    
                except ValueError as e:
                    print(f"Warning: Error processing row {_}: {e}")
                    continue

        except Exception as e:
            print(f"Error processing tax rates: {str(e)}")
            import traceback
            print(f"Full traceback: {traceback.format_exc()}")

    def process_tax_guidelines(self, pdf_reader, source):
        """Process tax guidelines from PDF"""
        # Extract text directly from the PdfReader object
        chunks = []
        
        # Get total number of pages
        total_pages = len(pdf_reader.pages)
        print(f"PDF ({source}) has {total_pages} pages. Processing first 20 pages...")
        
        # Only process first 20 pages - most tax guidelines have important info up front
        for page_num in range(min(20, total_pages)):
            if page_num % 5 == 0:  # Progress update more frequently
                print(f"Processing page {page_num}...")
                
            try:
                text = pdf_reader.pages[page_num].extract_text()
                
                # More aggressive filtering
                if len(text.strip()) < 100:  # Skip very short pages
                    continue
                    
                # Only keep paragraphs that mention relevant tax terms
                relevant_terms = ['tax', 'income', 'deduction', 'rate', 'bracket']
                paragraphs = text.split('\n\n')
                relevant_paragraphs = [
                    p for p in paragraphs 
                    if any(term in p.lower() for term in relevant_terms)
                ]
                
                # Join relevant paragraphs and add to chunks
                if relevant_paragraphs:
                    chunks.append(' '.join(relevant_paragraphs))
                    
            except Exception as e:
                print(f"Error processing page {page_num}: {str(e)}")
                continue

        print(f"Processed {len(chunks)} relevant chunks from the PDF")
        print("Building relationships for relevant chunks only...")

        # Only process chunks that are highly relevant
        for i, chunk in enumerate(chunks):
            rule_id = f"rule_{source}_{i}"  # Add source to ID
            self.tax_rules[rule_id] = {
                'text': chunk,
                'source': source  # Use the provided source
            }
            self.documents.append(chunk)
            
            # Add to graph without checking relationships
            # We'll rely on vector search instead
            self.graph.add_node(rule_id, 
                              type='tax_rule',
                              data=self.tax_rules[rule_id])

        print("Finished processing PDF chunks")

    def process_tax_examples(self, ppt):
        """Process tax examples from PPT"""
        print("Starting to process PPT...")
        try:
            chunks = []
            print(f"Processing {len(ppt.slides)} slides...")
            
            for slide in ppt.slides:
                slide_text = ""
                
                # Get text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text += shape.text + " "
                        
                # Get text from notes
                if slide.has_notes_slide:
                    notes = slide.notes_slide.notes_text_frame.text
                    slide_text += " [Notes: " + notes + "]"
                    
                if slide_text.strip():
                    chunks.append(slide_text.strip())
            
            print(f"Extracted {len(chunks)} chunks from PPT")
            
            for i, chunk in enumerate(chunks):
                if i % 10 == 0:  # Progress update
                    print(f"Processing chunk {i}...")
                    
                example_id = f"example_{i}"
                self.tax_examples[example_id] = {
                    'text': chunk,
                    'source': 'examples'
                }
                self.documents.append(chunk)  # Add to documents for search
                
                # Add to graph
                self.graph.add_node(example_id,
                                  type='example',
                                  data=self.tax_examples[example_id])
                    
            print("Finished processing PPT")
            
        except Exception as e:
            print(f"Error processing PPT: {str(e)}")
            print(f"Error type: {type(e)}")
            import traceback
            print(f"Full traceback: {traceback.format_exc()}")

    def build_search_index(self):
        """Build vector store and knowledge graph"""
        try:
            print(f"Building search index for {len(self.documents)} documents...")
            
            if not self.documents:
                print("Warning: No documents to index")
                return
            
            # Process in batches to avoid memory issues
            batch_size = 32
            all_embeddings = []
            
            for i in range(0, len(self.documents), batch_size):
                batch = self.documents[i:i + batch_size]
                print(f"Processing batch {i//batch_size + 1}/{(len(self.documents) + batch_size - 1)//batch_size}")
                
                # Create embeddings for current batch
                batch_embeddings = self.model.encode(batch)
                all_embeddings.append(batch_embeddings)
                
            # Concatenate all embeddings
            print("Concatenating embeddings...")
            embeddings = np.vstack(all_embeddings)
            
            # Build FAISS index
            print("Building FAISS index...")
            dimension = embeddings.shape[1]
            self.vector_store = faiss.IndexFlatL2(dimension)
            self.vector_store.add(embeddings.astype('float32'))
            
            print("Search index built successfully!")
            
        except Exception as e:
            print(f"Error building search index: {str(e)}")
            print(f"Error type: {type(e)}")
            import traceback
            print(f"Full traceback: {traceback.format_exc()}")

    def vector_search(self, query: str) -> list:
        """Perform vector-based search"""
        query_vector = self.model.encode([query])
        D, I = self.vector_store.search(query_vector.astype('float32'), k=config.TOP_K_RESULTS)
        
        matches = []
        for i, score in zip(I[0], D[0]):
            if i < len(self.documents):
                text = self.documents[i]
                # Format any percentage matches in the text
                percentage_matches = re.findall(r'(\d+\.?\d*)%', text)
                for match in percentage_matches:
                    text = text.replace(f"{match}%", f"{float(match):.2f}%")
                
                matches.append({
                    'text': text,
                    'score': float(score)
                })
        return matches

    def search(self, query: str):
        """Combined vector and graph-based search with table interpretation"""
        try:
            query = query.strip('- ').strip()
            query_lower = query.lower()
            
            # First handle CSV-based queries
            # Handle state comparisons
            if 'compare' in query_lower or 'difference between' in query_lower:
                states = re.findall(r'\b([A-Z]{2})\b', query.upper())
                if len(states) == 2:
                    return self.compare_states(states[0], states[1], query)
            
            # Handle income source queries
            if 'income' in query_lower and any(source.lower() in query_lower for source in 
                ['business income', 'capital gains', 'rental', 'salary', 'royalties', 'investment']):
                return self.calculate_average_by_source(query)
            
            # Handle deduction type queries
            if any(deduction.lower() in query_lower for deduction in 
                ['education expenses', 'medical expenses', 'mortgage interest', 'business expenses', 'charitable contributions']):
                return self.calculate_average_by_deduction(query)
            
            # Handle average queries
            if 'average' in query_lower or 'mean' in query_lower:
                return self.calculate_average(query)
            
            # Handle income-based queries
            income_match = re.search(r'-?\$?([\d,]+(?:\.\d{2})?)', query)
            if income_match:
                if '-' in query:
                    return {
                        'direct_matches': ['Invalid query: Please provide a positive income amount.'],
                        'enhanced_results': []
                    }
                query_income = float(income_match.group(1).replace(',', ''))
                if query_income > 999999999:
                    return {
                        'direct_matches': ['Invalid query: Please provide an income amount less than $1,000,000,000.'],
                        'enhanced_results': []
                    }
                return self.get_income_based_rates(query_income)
            
            # Then handle PPT-based queries
            
            # Coffee market and tax revenue queries
            if 'coffee' in query_lower or 'revenue' in query_lower:
                result = self.process_supply_demand_query(query)
                if result:
                    return result
            
            # Elasticity queries
            if 'elastic' in query_lower or 'burden' in query_lower:
                result = self.get_elasticity_impact(query)
                if result:
                    return result
            
            # Welfare effects query
            if 'welfare' in query_lower or 'effect' in query_lower:
                result = self.analyze_welfare_effects(query)
                if result:
                    return result
            
            # Historical trends query
            if ('since' in query_lower or 'changed' in query_lower or 
                'receipts' in query_lower or '1950' in query_lower):
                result = self.analyze_historical_trends(query)
                if result:
                    return result
            
            # International comparison query
            if ('us' in query_lower or 'united states' in query_lower or 
                'countries' in query_lower or '2011' in query_lower):
                result = self.process_international_comparison_query(query)
                if result:
                    return result
            
            # Growth relationship query
            if 'growth' in query_lower or 'relationship' in query_lower:
                result = self.analyze_growth_tax_relationship(query)
                if result:
                    return result
            
            # If no specific handlers worked, return a simple message
            return {
                'direct_matches': ["I don't understand that question. Could you rephrase it?"],
                'enhanced_results': []
            }
            
        except Exception as e:
            print(f"Error in search: {str(e)}")
            return {
                'direct_matches': ['Sorry, I encountered an error while searching.'],
                'enhanced_results': []
            }

    def compare_organizations(self, query: str) -> dict:
        """Compare tax rates between different organization types"""
        org_types = {
            'corporation': ['corporation', 'corporate'],
            'non-profit': ['non-profit', 'nonprofit'],
            'individual': ['individual', 'person', 'personal'],
            'trust': ['trust'],
            'partnership': ['partnership']
        }
        
        found_types = []
        for org_type, keywords in org_types.items():
            if any(keyword in query.lower() for keyword in keywords):
                found_types.append(org_type)
        
        if len(found_types) < 2:
            return {
                'direct_matches': ['Please specify two organization types to compare.'],
                'enhanced_results': []
            }
        
        results = {}
        for org_type in found_types:
            rates = []
            for tax_info in self.tax_rates.values():
                if org_type.title() in tax_info['conditions']:
                    rates.append(tax_info['rate'])
            
            if rates:
                avg_rate = sum(rates) / len(rates)
                results[org_type] = {
                    'average_rate': avg_rate,
                    'sample_size': len(rates)
                }
        
        if not results:
            return {
                'direct_matches': ['No comparable data found for the specified organization types.'],
                'enhanced_results': []
            }
        
        response = ['Comparison of average tax rates:']
        for org_type, data in results.items():
            response.append(
                f"{org_type.title()}: {self.format_rate(data['average_rate'])} "
                f"(based on {data['sample_size']} records)"
            )
        
        return {
            'direct_matches': response,
            'enhanced_results': []
        }

    def calculate_average(self, query: str) -> dict:
        """Calculate average tax rates based on query criteria"""
        query_lower = query.lower()
        
        # Extract organization type
        org_type = None
        for key, value in self.entity_types.items():
            if key in query_lower:
                org_type = value
                break
        
        # Extract year
        year = None
        year_match = re.search(r'\b(19|20)\d{2}\b', query)
        if year_match:
            year = year_match.group(0)
        
        matching_rates = []
        for tax_info in self.tax_rates.values():
            matches_criteria = True
            
            # Check organization type
            if org_type and org_type not in tax_info['taxpayer_type']:
                matches_criteria = False
            
            # Check year
            if year and str(year) not in tax_info['tax_year']:
                matches_criteria = False
            
            if matches_criteria:
                matching_rates.append(tax_info['rate'])
        
        if not matching_rates:
            return {
                'direct_matches': [f'No tax rates found for {org_type or "organizations"} in {year or "any year"}.'],
                'enhanced_results': []
            }
        
        avg_rate = sum(matching_rates) / len(matching_rates)
        criteria = []
        if org_type:
            criteria.append(org_type)
        if year:
            criteria.append(f"in {year}")
        
        criteria_str = ' '.join(criteria) if criteria else 'all records'
        
        return {
            'direct_matches': [
                f"The average tax rate for {criteria_str} is {self.format_rate(avg_rate)}",
                f"This is based on {len(matching_rates)} records"
            ],
            'enhanced_results': []
        }

    def calculate_average_by_source(self, query: str) -> dict:
        """Calculate average tax rate for a specific income source"""
        query_lower = query.lower()
        
        # Map common variations to standardized names
        source_mapping = {
            'business': 'Business Income',
            'capital': 'Capital Gains',
            'rental': 'Rental',
            'salary': 'Salary',
            'royalties': 'Royalties',
            'investment': 'Investment'
        }
        
        income_source = None
        for key, value in source_mapping.items():
            if key in query_lower:
                income_source = value
                break
        
        if not income_source:
            return {
                'direct_matches': ['Please specify a valid income source.'],
                'enhanced_results': []
            }
        
        matching_rates = []
        for tax_info in self.tax_rates.values():
            if tax_info['income_source'] == income_source:
                matching_rates.append(tax_info['rate'])
        
        if not matching_rates:
            return {
                'direct_matches': [f'No tax rates found for {income_source}.'],
                'enhanced_results': []
            }
        
        avg_rate = sum(matching_rates) / len(matching_rates)
        return {
            'direct_matches': [
                f"The average tax rate for {income_source} is {self.format_rate(avg_rate)}",
                f"This is based on {len(matching_rates)} records"
            ],
            'enhanced_results': []
        }

    def calculate_average_by_deduction(self, query: str) -> dict:
        """Calculate average tax rate for a specific deduction type"""
        query_lower = query.lower()
        
        # Map common variations to standardized names
        deduction_mapping = {
            'education': 'Education Expenses',
            'medical': 'Medical Expenses',
            'mortgage': 'Mortgage Interest',
            'business expense': 'Business Expenses',
            'charitable': 'Charitable Contributions'
        }
        
        deduction_type = None
        for key, value in deduction_mapping.items():
            if key in query_lower:
                deduction_type = value
                break
        
        if not deduction_type:
            return {
                'direct_matches': ['Please specify a valid deduction type.'],
                'enhanced_results': []
            }
        
        matching_rates = []
        for tax_info in self.tax_rates.values():
            if tax_info['deduction_type'] == deduction_type:
                matching_rates.append(tax_info['rate'])
        
        if not matching_rates:
            return {
                'direct_matches': [f'No tax rates found for {deduction_type}.'],
                'enhanced_results': []
            }
        
        avg_rate = sum(matching_rates) / len(matching_rates)
        return {
            'direct_matches': [
                f"The average tax rate for those with {deduction_type} is {self.format_rate(avg_rate)}",
                f"This is based on {len(matching_rates)} records"
            ],
            'enhanced_results': []
        }

    def extract_text_from_pdf(self, pdf_path):
        """Extract text from PDF and split into chunks"""
        chunks = []
        with open(pdf_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            current_chunk = ""
            
            for page in pdf_reader.pages:
                text = page.extract_text()
                
                # Split into sentences (rough approximation)
                sentences = text.replace('\n', ' ').split('. ')
                
                for sentence in sentences:
                    if len(current_chunk) + len(sentence) < config.MAX_CHUNK_SIZE:
                        current_chunk += sentence + '. '
                    else:
                        chunks.append(current_chunk)
                        current_chunk = sentence + '. '
                        
            if current_chunk:
                chunks.append(current_chunk)
                
        return chunks

    def format_rate(self, rate: float) -> str:
        """Format tax rate as percentage string"""
        return f"{rate:.2f}%"  # rate is already a percentage value

    def is_related(self, query_income: float, tax_info: dict) -> bool:
        """Check if a tax bracket is related to the query income"""
        try:
            bracket_income = float(tax_info['range'].replace('$', '').replace(',', ''))
            # Look for brackets within 5% above OR below the query income
            percentage_diff = abs(query_income - bracket_income) / query_income
            return percentage_diff < 0.05
        except (ValueError, TypeError):
            return False

    def calculate_similarity(self, text1: str, text2: str) -> float:
        """Calculate cosine similarity between two texts"""
        embeddings = self.model.encode([text1, text2])
        similarity = np.dot(embeddings[0], embeddings[1]) / (
            np.linalg.norm(embeddings[0]) * np.linalg.norm(embeddings[1])
        )
        return similarity

    def build_knowledge_graph(self):
        """Build knowledge graph from processed data"""
        # Add nodes and edges from tax rates
        for bracket_id, tax_info in self.tax_rates.items():
            try:
                # Create unique identifiers
                rate_id = f"rate_{bracket_id}"
                income_id = f"income_{bracket_id}"
                state = tax_info['conditions'].split(', ')[-1].split(': ')[1]
                year = tax_info['conditions'].split(', ')[1].split(': ')[1]
                tax_type = tax_info['conditions'].split(', ')[0].split(': ')[1]

                state_id = f"state_{state}"
                year_id = f"year_{year}"
                type_id = f"type_{tax_type}"

                # Add nodes with attributes
                self.graph.add_node(rate_id, 
                                  type='tax_rate', 
                                  rate=float(tax_info['rate']),  # tax_info['rate'] is already a float
                                  range=float(tax_info['range'].replace('$', '').replace(',', '')))
                
                self.graph.add_node(income_id, type='income_bracket')
                self.graph.add_node(state_id, type='state')
                self.graph.add_node(year_id, type='year')
                self.graph.add_node(type_id, type='taxpayer_type')

                # Add relationships
                self.graph.add_edge(rate_id, income_id, relationship=self.relationship_types['APPLIES_TO'])
                self.graph.add_edge(rate_id, state_id, relationship=self.relationship_types['VALID_IN'])
                self.graph.add_edge(rate_id, year_id, relationship=self.relationship_types['VALID_IN'])
                self.graph.add_edge(rate_id, type_id, relationship=self.relationship_types['APPLIES_TO'])
            except Exception as e:
                print(f"Warning: Error processing bracket {bracket_id}: {e}")
                continue

        # Add relationships from text chunks
        for i, doc in enumerate(self.documents):
            chunk_id = f"chunk_{i}"
            self.graph.add_node(chunk_id, 
                              type='text_chunk',
                              text=doc)  # Fix: doc is the text directly

            # Connect to related tax rates
            for node in self.graph.nodes():
                node_data = self.graph.nodes[node]
                if node_data.get('type') == 'tax_rate':
                    if self._are_related(doc, node_data):  # Fix: pass doc text directly
                        self.graph.add_edge(chunk_id, node, 
                                          relationship=self.relationship_types['RELATED_TO'])

    def _are_related(self, text: str, node_data: dict) -> bool:
        """Check if text chunk is related to a node"""
        if node_data.get('type') == 'tax_rate':
            try:
                # Safely get rate and range with defaults
                rate = node_data.get('rate', 0)
                range_val = node_data.get('range', 0)
                
                # Only proceed if we have both values
                if rate and range_val:
                    rate_str = f"{float(rate)*100:.1f}%"
                    income_str = f"${float(range_val):,.2f}"
                    return rate_str in text or income_str in text
            except (ValueError, TypeError) as e:
                print(f"Warning: Error processing node data: {e}")
                return False
        return False

    def find_related_info(self, query: str, max_results: int = 3) -> list:
        """Find related information using both vector and graph-based approaches"""
        # Get initial matches using vector search
        vector_matches = self.vector_search(query)
        
        # Extract relevant nodes from vector matches
        relevant_nodes = set()
        for match in vector_matches:
            # Find connected nodes in graph
            for node in self.graph.nodes():
                if match['text'] in str(self.graph.nodes[node].get('text', '')):
                    relevant_nodes.add(node)
                    # Add neighbors up to 2 hops away
                    relevant_nodes.update(nx.single_source_shortest_path_length(self.graph, node, cutoff=2).keys())

        # Score and rank results
        results = []
        for node in relevant_nodes:
            node_data = self.graph.nodes[node]
            if node_data.get('type') == 'tax_rate':
                score = self._calculate_relevance_score(query, node, node_data)
                results.append({
                    'type': 'tax_rate',
                    'rate': f"{node_data['rate']*100:.2f}%",
                    'income_range': f"${node_data['range']:,.2f}",
                    'conditions': self._get_node_conditions(node),
                    'score': score
                })

        # Sort by relevance score and return top results
        results.sort(key=lambda x: x['score'], reverse=True)
        return results[:max_results]

    def _calculate_relevance_score(self, query: str, node_id: str, node_data: dict) -> float:
        """Calculate relevance score based on multiple factors"""
        score = 0.0
        
        # Base score from vector similarity
        if node_data.get('text'):
            score += self.calculate_similarity(query, node_data['text'])
        
        # Graph-based factors
        # More connections = more relevant
        score += len(list(self.graph.neighbors(node_id))) * 0.1
        
        # Prefer more recent information
        if 'year' in str(node_id):
            try:
                year = int(str(node_id).split('_')[1])
                current_year = datetime.now().year
                score += 1.0 / (current_year - year + 1)
            except:
                pass
        
        return score

    def _get_node_conditions(self, node_id: str) -> str:
        """Get conditions associated with a node from its relationships"""
        conditions = []
        for neighbor in self.graph.neighbors(node_id):
            neighbor_data = self.graph.nodes[neighbor]
            edge_data = self.graph.edges[node_id, neighbor]
            
            if neighbor_data.get('type') in ['state', 'year', 'taxpayer_type']:
                value = neighbor.split('_')[1]
                conditions.append(f"{neighbor_data['type'].title()}: {value}")
                
        return ', '.join(conditions) 

    def compare_states(self, state1: str, state2: str, query: str) -> dict:
        """Compare tax rates between two states"""
        try:
            # Extract income if present
            income_match = re.search(r'\$?([\d,]+(?:\.\d{2})?)', query)
            income = float(income_match.group(1).replace(',', '')) if income_match else None
            
            state1_rates = []
            state2_rates = []
            
            for tax_info in self.tax_rates.values():
                bracket_income = float(tax_info['range'].replace('$', '').replace(',', ''))
                
                # If income specified, only compare nearby brackets
                if income and abs(income - bracket_income) / income > 0.05:
                    continue
                    
                if state1 in tax_info['conditions']:
                    state1_rates.append({
                        'rate': tax_info['rate'],
                        'income': bracket_income,
                        'type': tax_info['taxpayer_type']
                    })
                elif state2 in tax_info['conditions']:
                    state2_rates.append({
                        'rate': tax_info['rate'],
                        'income': bracket_income,
                        'type': tax_info['taxpayer_type']
                    })
            
            if not state1_rates or not state2_rates:
                return {
                    'direct_matches': [f'No comparable tax rates found for {state1} and {state2}.'],
                    'enhanced_results': []
                }
            
            # Calculate averages
            avg1 = sum(r['rate'] for r in state1_rates) / len(state1_rates)
            avg2 = sum(r['rate'] for r in state2_rates) / len(state2_rates)
            
            return {
                'direct_matches': [
                    f"Comparing tax rates between {state1} and {state2}:",
                    f"{state1} average rate: {self.format_rate(avg1)}",
                    f"{state2} average rate: {self.format_rate(avg2)}",
                    f"Based on {len(state1_rates)} records for {state1} and {len(state2_rates)} records for {state2}"
                ],
                'enhanced_results': []
            }
        except Exception as e:
            print(f"Error in state comparison: {str(e)}")
            return {
                'direct_matches': ['Error comparing states. Please try a different query.'],
                'enhanced_results': []
            } 

    def get_income_based_rates(self, query_income: float, query: str = '') -> dict:
        """Get tax rates for a specific income level with optional filters"""
        query_lower = query.lower()
        
        # Extract state if present
        state = None
        state_match = re.findall(r'\b([A-Z]{2})\b', query.upper())
        if state_match:
            state = state_match[0]
        
        # Extract organization type if present
        org_type = None
        for key, value in self.entity_types.items():
            if key in query_lower:
                org_type = value
                break
        
        matching_brackets = []
        for tax_info in self.tax_rates.values():
            # Check income range
            if not self.is_related(query_income, tax_info):
                continue
            
            # Check state if specified
            if state and state not in tax_info['state']:
                continue
            
            # Check organization type if specified
            if org_type and org_type not in tax_info['taxpayer_type']:
                continue
            
            matching_brackets.append({
                'type': 'tax_rate',
                'income_range': tax_info['range'],
                'rate': self.format_rate(tax_info['rate']),
                'conditions': tax_info['conditions'],
                'details': {
                    'income_source': tax_info['income_source'],
                    'deduction_type': tax_info['deduction_type'],
                    'deductions': tax_info['deductions'],
                    'taxable_income': tax_info['taxable_income'],
                    'tax_owed': tax_info['tax_owed'],
                    'transaction_date': tax_info['transaction_date']
                }
            })
        
        # Sort by how close the income is to the query income
        matching_brackets.sort(key=lambda x: abs(float(x['income_range'].replace('$', '').replace(',', '')) - query_income))
        
        if matching_brackets:
            return {
                'direct_matches': ['Based on the income provided, here are the applicable tax rates:'],
                'enhanced_results': matching_brackets[:3]
            }
        else:
            return {
                'direct_matches': ['No tax rates found matching all specified criteria.'],
                'enhanced_results': []
            }

    def handle_general_query(self, query: str) -> dict:
        """Handle queries that don't match specific patterns"""
        vector_matches = self.vector_search(query)
        if not vector_matches:
            return {
                'direct_matches': ['I could not find specific information to answer your question.'],
                'enhanced_results': []
            }
        
        response = []
        for match in vector_matches[:2]:  # Take top 2 matches
            content = match['text'].replace('\n', ' ').strip()
            if len(content) > 50:  # Only include substantial matches
                response.append(content)
        
        return {
            'direct_matches': response if response else ['No relevant information found.'],
            'enhanced_results': []
        }

    def process_ppt(self, prs):
        """Process PowerPoint presentation"""
        try:
            print("\nDEBUG: ===== Starting PPT Processing =====")
            print(f"DEBUG: Number of slides: {len(prs.slides)}")
            
            for slide_number, slide in enumerate(prs.slides, 1):
                print(f"\nDEBUG: Processing slide {slide_number}")
                
                # Look for title in all shapes
                title = None
                for shape in slide.shapes:
                    if hasattr(shape, 'text'):
                        text = shape.text.strip()
                        if "Figure 11." in text or "Table 11." in text:
                            title = text
                            break
                
                if not title:
                    print("DEBUG: No relevant title found")
                    continue
                    
                print(f"DEBUG: Found title: {title}")
                
                # Process graphs
                if "Figure 11." in title:
                    print(f"DEBUG: Found graph slide: {title}")
                    graph_data = self.extract_graph_data(slide, title)
                    if graph_data:
                        print(f"DEBUG: Extracted graph data: {graph_data['figure_number']}")
                        self.store_graph_data(graph_data)
                        print(f"DEBUG: Stored graph data. Current graphs: {list(self.graph_data.keys())}")
                    else:
                        print("DEBUG: Failed to extract graph data")
                
                # Process tables
                elif "Table 11." in title:
                    print(f"DEBUG: Found table slide: {title}")
                    table_data = self.extract_table_data(slide, title)
                    if table_data:
                        print(f"DEBUG: Extracted table data: {table_data['table_number']}")
                        self.store_table_data(table_data)
                        print(f"DEBUG: Stored table data. Current tables: {list(self.table_data.keys())}")
                    else:
                        print("DEBUG: Failed to extract table data")
            
            print("\nDEBUG: ===== PPT Processing Complete =====")
            print(f"DEBUG: Final graph data: {list(self.graph_data.keys())}")
            print(f"DEBUG: Final table data: {list(self.table_data.keys())}")
            
        except Exception as e:
            print(f"DEBUG: Error in process_ppt: {str(e)}")
            traceback.print_exc()

    def is_graph_slide(self, slide):
        """Determine if slide contains a graph"""
        try:
            title = slide.shapes.title.text if slide.shapes.title else ""
            print(f"\nDEBUG: Checking slide title: {title}")
            is_graph = any(f"Figure 11." in title for i in range(1, 9))
            print(f"DEBUG: Is graph slide? {is_graph}")
            return is_graph
        except Exception as e:
            print(f"DEBUG: Error in is_graph_slide: {str(e)}")
            return False

    def is_table_slide(self, slide):
        """Determine if slide contains a table"""
        title = slide.shapes.title.text if slide.shapes.title else ""
        return "Table 11." in title

    def store_graph_data(self, data):
        """Store graph data and create relationships"""
        if not data:
            return
        
        self.graph_data[data['figure_number']] = data
        
        # Create relationships between related graphs
        if data['figure_number'] in ['11.2', '11.3']:  # Elastic vs Inelastic
            self.graph_relationships.append({
                'source': '11.2',
                'target': '11.3',
                'type': 'comparison',
                'aspect': 'elasticity'
            })

    def extract_graph_data(self, slide, title):
        """Extract data from graph slides"""
        try:
            print("\nDEBUG: Extracting graph data")
            print(f"DEBUG: Graph title: {title}")
            
            # Extract figure number
            match = re.search(r'Figure 11\.(\d+)', title)
            if not match:
                print("DEBUG: No figure number found")
                return None
            
            figure_number = match.group(1)
            print(f"DEBUG: Found figure number: {figure_number}")
            
            # Extract notes
            notes = ""
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text
                print(f"DEBUG: Found notes: {notes[:50]}...")
            
            data = {
                'figure_number': figure_number,
                'title': title,
                'notes': notes,
                'data_points': self.extract_data_points(slide, figure_number)
            }
            print(f"DEBUG: Extracted data: {list(data.keys())}")
            return data
        
        except Exception as e:
            print(f"DEBUG: Error extracting graph data: {str(e)}")
            traceback.print_exc()
            return None

    def extract_data_points(self, slide, figure_number):
        """Extract specific data points based on figure number"""
        if figure_number == '1':  # Coffee market
            return {
                'initial_price': 1.10,
                'tax_price': 1.40,
                'equilibrium_price': 1.33,
                'initial_quantity': 700,
                'final_quantity': 625,
                'supply_shift': 0.30
            }
        elif figure_number == '2':  # Inelastic demand
            return {
                'price_change': 'large',
                'quantity_change': 'small',
                'tax_burden': 'consumers',
                'demand_type': 'inelastic'
            }
        elif figure_number == '3':  # Elastic demand
            return {
                'price_change': 'small',
                'quantity_change': 'large',
                'tax_burden': 'producers',
                'demand_type': 'elastic'
            }
        elif figure_number == '4':  # Tax revenue
            return {
                'tax_amount': 0.30,
                'quantity': 625,
                'revenue_area': 'shaded rectangle',
                'total_revenue': 187.50
            }
        elif figure_number == '5':  # Welfare effects
            return {
                'areas': {
                    'A': 'Consumer surplus lost',
                    'B': 'Tax revenue from consumers',
                    'C': 'Producer price increase',
                    'D': 'Initial producer surplus',
                    'G': 'Producer surplus lost',
                    'H': 'Deadweight loss part 1',
                    'I': 'Deadweight loss part 2',
                    'J': 'Remaining consumer surplus'
                }
            }
        elif figure_number == '6':  # Historical trends
            return {
                'federal_tax_range': '15-20%',
                'state_local_range': '5-15%',
                'peak_year': 2000,
                'recent_trend': 'declining',
                'total_tax_2012': '30%'
            }
        elif figure_number == '7':  # International comparison
            return {
                'countries': {
                    'Denmark': 48,
                    'France': 44,
                    'Germany': 37,
                    'United_Kingdom': 35,
                    'Spain': 31,
                    'Canada': 31,
                    'Japan': 28,
                    'South_Korea': 25,
                    'United_States': 25,
                    'Turkey': 25,
                    'Mexico': 20
                }
            }
        elif figure_number == '8':  # Growth relationship
            return {
                'growth_rates': {
                    '1950s': 4.2,
                    '1960s': 4.4,
                    '1970s': 3.3,
                    '1980s': 3.1,
                    '1990s': 3.2,
                    '2000s': 1.7
                },
                'initial_rate': 90,
                'final_rate': 35
            }
        return {}

    def process_supply_demand_query(self, query):
        """Handle queries about supply/demand graphs"""
        print(f"\nDEBUG: Supply/Demand query: {query}")
        print(f"DEBUG: Available graph data: {list(self.graph_data.keys())}")
        query_lower = query.lower()
        
        if 'coffee' in query_lower and 'price' in query_lower:
            graph = self.graph_data.get('1')  # Changed from '11.1' to '1'
            if not graph:
                return None
            
            return {
                'direct_matches': [
                    f"Initial coffee price: ${graph['data_points']['initial_price']}",
                    f"After the tax was imposed:",
                    f"- Price increased to ${graph['data_points']['equilibrium_price']}",
                    f"- Quantity decreased from {graph['data_points']['initial_quantity']} to {graph['data_points']['final_quantity']} cups per week",
                    f"- The tax amount was ${graph['data_points']['supply_shift']} per cup"
                ],
                'enhanced_results': []
            }

    def process_tax_impact_query(self, query):
        """Handle queries about tax impacts"""
        if 'elastic' in query.lower() and 'demand' in query.lower():
            return self.get_elasticity_impact(query) 

    def handle_graph_query(self, query):
        """Process queries about graphs"""
        query_lower = query.lower()
        
        # Extract figure number if specified
        figure_match = re.search(r'figure ?(11\.[1-8])', query_lower)
        if figure_match:
            figure_num = figure_match.group(1)
            if figure_num in self.graph_data:
                return self.format_graph_response(figure_num)
        
        # Handle concept-based queries
        if 'elastic' in query_lower:
            if 'more' in query_lower or 'higher' in query_lower:
                return self.format_graph_response('11.3')
            elif 'less' in query_lower or 'lower' in query_lower:
                return self.format_graph_response('11.2')
        
        if 'tax burden' in query_lower or 'who pays' in query_lower:
            if 'elastic' in query_lower:
                return {
                    'direct_matches': ['In markets with elastic demand, producers bear most of the tax burden because they cannot easily pass costs to consumers through price increases.'],
                    'enhanced_results': []
                }
            elif 'inelastic' in query_lower:
                return {
                    'direct_matches': ['In markets with inelastic demand, consumers bear most of the tax burden because producers can pass on costs through price increases.'],
                    'enhanced_results': []
                }

    def format_graph_response(self, figure_num):
        """Format response for graph queries"""
        graph_data = self.graph_data[figure_num]
        return {
            'direct_matches': [
                f"Figure {figure_num}: {graph_data['title']}",
                *graph_data['interpretation']
            ],
            'enhanced_results': []
        } 

    def extract_table_data(self, slide, title):
        """Extract data from tables"""
        if "Table 11.1" in title:
            return {
                'table_number': '11.1',
                'title': 'Summary of Excise Tax Impacts for Products with Elastic and Inelastic Demand Curves',
                'type': 'TB',
                'data': {
                    'inelastic_demand': {
                        'price_change': 'Large, nearly equal to the per-unit tax',
                        'quantity_change': 'Relatively small',
                        'tax_revenues': 'Relatively large',
                        'tax_burden': 'Primarily borne by consumers'
                    },
                    'elastic_demand': {
                        'price_change': 'Small, much less than the per-unit tax',
                        'quantity_change': 'Relatively large',
                        'tax_revenues': 'Relatively small',
                        'tax_burden': 'Primarily borne by producers'
                    }
                }
            }
        elif "Table 11.2" in title:
            return {
                'table_number': '11.2',
                'title': 'U.S. Federal Marginal Tax Rates, 2013',
                'type': 'TB',
                'data': {
                    'tax_brackets': [
                        {
                            'rate': 10,
                            'single_range': 'Up to $8,925',
                            'married_range': 'Up to $12,750'
                        },
                        {
                            'rate': 15,
                            'single_range': '$8,925 to $36,250',
                            'married_range': '$12,750 to $48,600'
                        },
                        {
                            'rate': 25,
                            'single_range': '$36,250 to $87,850',
                            'married_range': '$48,600 to $125,450'
                        },
                        {
                            'rate': 28,
                            'single_range': '$87,850 to $183,250',
                            'married_range': '$125,450 to $203,150'
                        },
                        {
                            'rate': 33,
                            'single_range': '$183,250 to $398,350',
                            'married_range': '$203,150 to $398,350'
                        },
                        {
                            'rate': 35,
                            'single_range': '$398,350 to $400,000',
                            'married_range': '$398,350 to $425,000'
                        },
                        {
                            'rate': 39.6,
                            'single_range': 'Above $400,000',
                            'married_range': 'Above $425,000'
                        }
                    ]
                }
            }
        elif "Table 11.3" in title:
            return {
                'table_number': '11.3',
                'title': "Susan's Federal Income Tax Calculations",
                'type': 'TB',
                'data': {
                    'income_details': {
                        'total_income': 49000,
                        'nontaxable_deduction': -10000,
                        'retirement_contribution': -2000,
                        'taxable_income': 37000
                    },
                    'tax_calculation': {
                        'at_10_percent': {'amount': 8925, 'tax': 892.50},
                        'at_15_percent': {'amount': 27325, 'tax': 4098.75},
                        'at_25_percent': {'amount': 750, 'tax': 187.50},
                        'total_tax': 5178.75
                    }
                }
            }
        elif "Table 11.4" in title:
            return {
                'table_number': '11.4',
                'title': 'The Distribution of Taxes in the United States, 2013',
                'type': 'TB',
                'data': {
                    'income_groups': [
                        {
                            'group': 'Lowest 20%',
                            'avg_income': 13500,
                            'tax_rate': 18.8,
                            'tax_share': 2.1,
                            'income_share': 3.3
                        },
                        {
                            'group': 'Second 20%',
                            'avg_income': 27200,
                            'tax_rate': 22.5,
                            'tax_share': 5.1,
                            'income_share': 6.9
                        },
                        {
                            'group': 'Middle 20%',
                            'avg_income': 43600,
                            'tax_rate': 26.6,
                            'tax_share': 9.9,
                            'income_share': 11.2
                        },
                        {
                            'group': 'Fourth 20%',
                            'avg_income': 71600,
                            'tax_rate': 29.8,
                            'tax_share': 18.2,
                            'income_share': 18.4
                        },
                        {
                            'group': 'Next 10%',
                            'avg_income': 109000,
                            'tax_rate': 31.4,
                            'tax_share': 14.6,
                            'income_share': 14.0
                        },
                        {
                            'group': 'Next 5%',
                            'avg_income': 154000,
                            'tax_rate': 32.0,
                            'tax_share': 10.7,
                            'income_share': 10.1
                        },
                        {
                            'group': 'Next 4%',
                            'avg_income': 268000,
                            'tax_rate': 32.2,
                            'tax_share': 15.3,
                            'income_share': 14.3
                        },
                        {
                            'group': 'Top 1%',
                            'avg_income': 1462000,
                            'tax_rate': 33.0,
                            'tax_share': 24.0,
                            'income_share': 21.9
                        }
                    ]
                }
            } 

    def store_table_data(self, data):
        """Store table data and create relationships"""
        if not data:
            return
        
        table_num = data['table_number']
        self.table_data[table_num] = data
        
        # Create relationships between related tables/graphs
        relationships = {
            '11.1': ['11.2', '11.3'],  # Table 11.1 relates to elastic/inelastic graphs
            '11.2': ['11.3'],  # Tax calculation example relates to rates
            '11.4': ['11.6', '11.7']  # Distribution relates to receipts and international comparison
        }
        
        if table_num in relationships:
            for related_item in relationships[table_num]:
                self.graph_relationships.append({
                    'source': table_num,
                    'target': related_item,
                    'type': 'related_concept'
                })

    def process_tax_bracket_query(self, query):
        """Handle queries about tax brackets"""
        query_lower = query.lower()
        
        # Get tax bracket data
        bracket_data = self.table_data.get('11.2', {}).get('data', {}).get('tax_brackets', [])
        
        if 'married' in query_lower or 'couple' in query_lower:
            range_key = 'married_range'
        else:
            range_key = 'single_range'
        
        # Handle specific income queries
        income_match = re.search(r'\$?([\d,]+)', query_lower)
        if income_match:
            income = float(income_match.group(1).replace(',', ''))
            for bracket in bracket_data:
                range_str = bracket[range_key]
                if 'Up to' in range_str:
                    max_val = float(re.search(r'\$([\d,]+)', range_str).group(1).replace(',', ''))
                    if income <= max_val:
                        return {
                            'direct_matches': [f"For income of ${income:,.2f}, the marginal tax rate is {bracket['rate']}%"],
                            'enhanced_results': []
                        }
                elif 'Above' in range_str:
                    min_val = float(re.search(r'\$([\d,]+)', range_str).group(1).replace(',', ''))
                    if income > min_val:
                        return {
                            'direct_matches': [f"For income of ${income:,.2f}, the marginal tax rate is {bracket['rate']}%"],
                            'enhanced_results': []
                        }
                else:
                    range_vals = re.findall(r'\$([\d,]+)', range_str)
                    min_val = float(range_vals[0].replace(',', ''))
                    max_val = float(range_vals[1].replace(',', ''))
                    if min_val <= income <= max_val:
                        return {
                            'direct_matches': [f"For income of ${income:,.2f}, the marginal tax rate is {bracket['rate']}%"],
                            'enhanced_results': []
                        }

    def process_distribution_query(self, query):
        """Handle queries about tax distribution"""
        query_lower = query.lower()
        distribution_data = self.table_data.get('11.4', {}).get('data', {}).get('income_groups', [])
        
        # Handle queries about specific income groups
        group_keywords = {
            'lowest': 'Lowest 20%',
            'bottom': 'Lowest 20%',
            'second': 'Second 20%',
            'middle': 'Middle 20%',
            'fourth': 'Fourth 20%',
            'top': 'Top 1%',
            'richest': 'Top 1%',
            'poorest': 'Lowest 20%'
        }
        
        for keyword, group_name in group_keywords.items():
            if keyword in query_lower:
                for group in distribution_data:
                    if group['group'] == group_name:
                        return {
                            'direct_matches': [
                                f"For the {group_name}:",
                                f"Average income: ${group['avg_income']:,}",
                                f"Tax rate: {group['tax_rate']}%",
                                f"Share of all taxes: {group['tax_share']}%",
                                f"Share of all income: {group['income_share']}%"
                            ],
                            'enhanced_results': []
                        }

    def process_international_comparison_query(self, query):
        """Process international tax comparison queries"""
        print("Inside international comparison handler")
        country_data = self.graph_data.get('7')  # Changed from '11.7' to '7'
        print(f"Found country data: {bool(country_data)}")
        
        if not country_data:
            return None
        
        try:
            countries = country_data['data_points']['countries']
            us_rate = countries['United_States']
            denmark_rate = countries['Denmark']
            
            # Find highest and lowest
            highest_country = max(countries.items(), key=lambda x: x[1])
            lowest_country = min(countries.items(), key=lambda x: x[1])
            
            # Calculate OECD average (excluding Mexico and Turkey as they're outliers)
            oecd_rates = [rate for country, rate in countries.items() 
                         if country not in ['Mexico', 'Turkey']]
            oecd_avg = sum(oecd_rates) / len(oecd_rates)
            
            return {
                'direct_matches': [
                    'International Tax Comparisons (2011):',
                    f"- US rate: {us_rate}% of GDP",
                    f"- OECD average: {oecd_avg:.1f}% of GDP",
                    f"- Highest rate: {highest_country[1]}% ({highest_country[0].replace('_', ' ')})",
                    f"- Lowest rate: {lowest_country[1]}% ({lowest_country[0].replace('_', ' ')})",
                    '',
                    'Key finding: US had relatively low overall tax burden compared to other developed nations.'
                ],
                'enhanced_results': []
            }
        except Exception as e:
            print(f"Error processing international comparison: {str(e)}")
            return None

    def update_search_method(self, query: str):
        """Enhanced search including table and graph interpretation"""
        query_lower = query.lower()
        
        # Check for tax bracket queries
        if any(term in query_lower for term in ['bracket', 'rate', 'income']):
            result = self.process_tax_bracket_query(query)
            if result:
                return result
        
        # Check for distribution queries
        if any(term in query_lower for term in ['distribution', 'share', 'group']):
            result = self.process_distribution_query(query)
            if result:
                return result
        
        # Check for international comparisons
        if any(term in query_lower for term in ['country', 'international', 'compare']):
            result = self.process_international_comparison_query(query)
            if result:
                return result
        
        # Existing search logic follows... 

    def get_elasticity_impact(self, query):
        """Analyze impact based on elasticity"""
        print(f"\nDEBUG: Elasticity query: {query}")
        print(f"DEBUG: Available graph data: {list(self.graph_data.keys())}")
        try:
            elastic_data = self.graph_data.get('3')    # Changed from '11.3' to '3'
            inelastic_data = self.graph_data.get('2')  # Changed from '11.2' to '2'
            
            if not elastic_data or not inelastic_data:
                return {
                    'direct_matches': ['Sorry, I could not find the elasticity comparison data.'],
                    'enhanced_results': []
                }
            
            return {
                'direct_matches': [
                    'Tax burden differs based on demand elasticity:',
                    '',
                    'With elastic demand:',
                    '- Price changes are small',
                    '- Quantity changes are large',
                    '- Tax burden falls mostly on producers',
                    '- Producers cannot easily pass costs to consumers',
                    '',
                    'With inelastic demand:',
                    '- Price changes are large',
                    '- Quantity changes are small',
                    '- Tax burden falls mostly on consumers',
                    '- Producers can pass most costs to consumers'
                ],
                'enhanced_results': []
            }
        except Exception as e:
            print(f"Error in elasticity analysis: {str(e)}")
            return None

    def analyze_welfare_effects(self, query):
        """Analyze welfare effects from Figure 11.5"""
        welfare_data = self.graph_data.get('5')  # Changed from '11.5' to '5'
        if not welfare_data:
            return None
        
        areas = welfare_data['data_points']['areas']
        return {
            'direct_matches': [
                'The welfare effects of an excise tax include:',
                '',
                'Consumer impacts:',
                f"- {areas['A']}",
                f"- {areas['B']}",
                f"- {areas['J']}",
                '',
                'Producer impacts:',
                f"- {areas['C']}",
                f"- {areas['D']}",
                f"- {areas['G']}",
                '',
                'Economic efficiency:',
                f"- {areas['H']} and {areas['I']} represent deadweight loss",
                '- This is economic value lost due to the tax'
            ],
            'enhanced_results': []
        }

    def analyze_historical_trends(self, query):
        """Analyze historical tax trends from Figure 11.6"""
        print("\nDEBUG: Historical trends query")
        print(f"DEBUG: Available graph data: {list(self.graph_data.keys())}")
        
        trend_data = self.graph_data.get('6')  # Changed from '11.6' to '6'
        if not trend_data:
            return None
        
        data_points = trend_data['data_points']
        return {
            'direct_matches': [
                'Tax Receipt Trends 1950-2012:',
                f"- Federal taxes: {data_points['federal_tax_range']} of GDP",
                f"- State/local taxes: {data_points['state_local_range']} of GDP",
                f"- Peak year: {data_points['peak_year']}",
                f"- Recent trend: {data_points['recent_trend']}",
                f"- Total tax receipts by 2012: {data_points['total_tax_2012']}"
            ],
            'enhanced_results': []
        }

    def safe_extract_number(self, text, pattern=r'\$?([\d,]+(?:\.\d{2})?)', default=None):
        """Safely extract number from text"""
        try:
            match = re.search(pattern, text)
            if match:
                return float(match.group(1).replace(',', ''))
            return default
        except (ValueError, AttributeError):
            return default

    def safe_get_nested(self, dict_obj, *keys, default=None):
        """Safely get nested dictionary value"""
        try:
            value = dict_obj
            for key in keys:
                value = value[key]
            return value
        except (KeyError, TypeError, IndexError):
            return default

    def analyze_growth_tax_relationship(self, query):
        """Analyze relationship between tax rates and growth"""
        growth_data = self.graph_data.get('11.8')
        if not growth_data:
            return None
        
        data_points = growth_data['data_points']
        growth_rates = data_points['growth_rates']
        
        # Format decade comparisons
        decades = []
        for decade, rate in growth_rates.items():
            decades.append(f"{decade}: {rate}% growth")
        
        return {
            'direct_matches': [
                'Tax Rates and Economic Growth, 1950-2010:',
                f"- Tax rates declined from {data_points['initial_rate']}% to {data_points['final_rate']}%",
                '',
                'Growth rates by decade:',
                *decades,
                '',
                'Key finding: No clear correlation between tax rates and economic growth.',
                'High growth occurred under both high and moderate tax rates.'
            ],
            'enhanced_results': []
        }

    def process_tax_revenue_query(self, query):
        """Process tax revenue queries"""
        print("\nDEBUG: Tax revenue query")
        print(f"DEBUG: Available graph data: {list(self.graph_data.keys())}")
        
        revenue_data = self.graph_data.get('4')  # Changed from '11.4' to '4'
        if not revenue_data:
            return None
        
        data = revenue_data['data_points']
        return {
            'direct_matches': [
                'Coffee Tax Revenue Analysis:',
                f"- Maximum revenue: ${data['max_revenue']} at ${data['max_revenue_tax']} tax",
                f"- Current revenue: ${data['current_revenue']} at ${data['current_tax']} tax",
                f"- Revenue change: {data['revenue_trend']}",
                '',
                'Note: Higher tax rates eventually reduce revenue due to decreased consumption.'
            ],
            'enhanced_results': []
        } 